import os
from typing import List
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


class DocumentProcessor:
    """Process different document types and extract text"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PDF: {str(e)}")
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {str(e)}")
    
    @staticmethod
    def process_document(file_path: str, file_type: str) -> str:
        """Process document based on file type"""
        if file_type == "pdf":
            return DocumentProcessor.extract_text_from_pdf(file_path)
        elif file_type == "docx":
            return DocumentProcessor.extract_text_from_docx(file_path)
        elif file_type == "pptx":
            return DocumentProcessor.extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split text into chunks with overlap for better context retention
        
        Args:
            text: The text to chunk
            chunk_size: Maximum size of each chunk in characters
            overlap: Number of characters to overlap between chunks
            
        Returns:
            List of text chunks
        """
        if not text or len(text) == 0:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            # Calculate end position
            end = start + chunk_size
            
            # If this is not the last chunk, try to break at a sentence or word boundary
            if end < text_length:
                # Look for sentence boundary (. ! ?)
                sentence_end = max(
                    text.rfind('. ', start, end),
                    text.rfind('! ', start, end),
                    text.rfind('? ', start, end)
                )
                
                if sentence_end != -1 and sentence_end > start + chunk_size // 2:
                    end = sentence_end + 1
                else:
                    # Look for word boundary (space)
                    space_pos = text.rfind(' ', start, end)
                    if space_pos != -1 and space_pos > start + chunk_size // 2:
                        end = space_pos
            
            # Extract chunk
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - overlap if end < text_length else text_length
        
        return chunks
